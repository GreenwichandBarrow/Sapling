from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = "/tmp/sidney_garber_growth_section_mockup_2026-07-30.pptx"

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(118, 118, 118)
LIGHT_GRAY = RGBColor(240, 240, 240)
SIDEBAR = RGBColor(248, 241, 237)
PEACH = RGBColor(251, 224, 207)
PEACH2 = RGBColor(252, 237, 228)
ORANGE = RGBColor(189, 103, 35)
LIGHT_ORANGE = RGBColor(250, 190, 136)
LINE = RGBColor(199, 139, 95)


def add_text(slide, text, x, y, w, h, size=18, bold=False, italic=False,
             color=BLACK, align=None, font="Arial"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    r = p.runs[0]
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_bullets(slide, bullets, x, y, w, h, size=13.2, color=BLACK,
                leading=1.0, font="Arial"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + b
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6 * leading)
    return box


def add_rect(slide, x, y, w, h, fill=PEACH, line=LINE, width=1.0, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(width)
    return s


def add_header(slide, title, subtitle, page, source="Source: Company materials; management discussions; investor discussions; diligence to confirm"):
    # left band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.64), Inches(9))
    band.fill.solid()
    band.fill.fore_color.rgb = SIDEBAR
    band.line.fill.background()
    add_text(slide, "G\nB", 0.16, 0.25, 0.28, 0.55, size=17, color=BLACK, align=1)
    v = add_text(slide, "Luxury Jewelry Industry", 0.12, 5.6, 0.42, 2.25, size=16, color=GRAY, align=1)
    v.rotation = 270

    add_text(slide, title, 1.18, 0.55, 10.8, 0.45, size=25, color=BLACK)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(1.12), Inches(13.2), Inches(0.012))
    rule.fill.solid()
    rule.fill.fore_color.rgb = LINE
    rule.line.fill.background()
    add_text(slide, "SIDNEY\nGARBER", 13.25, 0.45, 1.5, 0.62, size=15, bold=True, align=1)
    add_text(slide, subtitle, 1.18, 1.28, 13.4, 0.55, size=15.5, italic=True, color=BLACK)

    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(8.62), Inches(14.0), Inches(0.015))
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(184, 184, 184)
    foot.line.fill.background()
    add_text(slide, "GB", 1.05, 8.70, 0.3, 0.18, size=8, bold=True, color=ORANGE)
    add_text(slide, source, 1.52, 8.70, 8.7, 0.18, size=6.8, color=GRAY)
    conf = add_text(slide, "Strictly Confidential", 15.23, 6.75, 0.28, 1.45, size=8, color=BLACK)
    conf.rotation = 270
    add_text(slide, str(page), 14.25, 8.38, 0.32, 0.18, size=8, color=GRAY)


def add_arrow_panel(slide, title, bullets, x, y, w, h, fill=RGBColor(235, 244, 255), line=RGBColor(47, 77, 111)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill
    arrow.line.color.rgb = line
    arrow.line.width = Pt(1.4)
    add_text(slide, title, x + 0.22, y + 0.10, w - 0.8, 0.22, size=13.5, bold=True, color=ORANGE)
    add_bullets(slide, bullets, x + 0.25, y + 0.36, w - 0.8, h - 0.38, size=9.2)
    return arrow


def slide_growth_strategy():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Multi-Pronged Growth Strategy",
        "Growth should be framed as diligence-backed controllable levers, not a forecast until seller financials and channel data are received",
        66,
    )

    # Left segmented concept using stacked wedges as a visual cue to the reference without recreating it exactly
    hub = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(0.95), Inches(2.55), Inches(2.05), Inches(3.75))
    hub.line.color.rgb = RGBColor(225, 225, 225)
    hub.line.width = Pt(26)
    add_text(slide, "SIDNEY\nGARBER", 0.95, 4.05, 1.5, 0.45, size=16, bold=True, color=ORANGE, align=1)

    labels = [
        ("Distribution\nExpansion", 0.47, 2.35, RGBColor(189, 103, 35)),
        ("Clienteling &\nPrivate Client", 1.30, 2.90, RGBColor(211, 130, 63)),
        ("Inventory\nProductivity", 1.55, 3.78, RGBColor(229, 158, 93)),
        ("Service Layer /\nRenewal", 1.42, 4.72, RGBColor(242, 191, 143)),
        ("Category /\nProduct Extension", 0.64, 5.40, RGBColor(246, 218, 196)),
    ]
    for txt, x, y, col in labels:
        s = slide.shapes.add_shape(MSO_SHAPE.PIE, Inches(x), Inches(y), Inches(1.35), Inches(0.72))
        s.fill.solid()
        s.fill.fore_color.rgb = col
        s.line.fill.background()
        add_text(slide, txt, x + 0.08, y + 0.16, 1.08, 0.4, size=9.5, bold=True,
                 italic=True, color=WHITE if col != RGBColor(246, 218, 196) else BLACK, align=1)

    rows = [
        ("Distribution Expansion", [
            "Expand selective wholesale / stockist relationships where brand fit and economics are attractive",
            "Improve ecommerce discovery and conversion without diluting high-touch positioning",
            "Test high-affluent markets through partners, trunk shows, and private-client events",
        ]),
        ("Clienteling & Private Client Activation", [
            "Institutionalize CRM, occasion-based outreach, and repeat-client cadence",
            "Convert relationship knowledge into measurable appointment, event, and repeat-purchase activity",
            "Reduce dependence on founder memory by building structured client data",
        ]),
        ("Inventory Productivity", [
            "Build SKU-level visibility into aging, turns, replenishment, markdowns, and margin by channel",
            "Prioritize working-capital efficiency without impairing brand presentation",
            "Use inventory data to separate collectible value from slow-moving stock",
        ]),
        ("Service Layer / Jewelry Renewal", [
            "Formalize repairs, care, repurposing, and bespoke refresh into repeatable service touchpoints",
            "Use service moments to deepen trust, client data, and future purchase intent",
            "Test whether client care can create less cyclical recurring revenue",
        ]),
        ("Category / Product Extension", [
            "Explore watches, everyday fine jewelry, bridal / milestone gifting, collaborations, and limited collections where brand permission exists",
            "Use small tests before committing inventory capital",
            "Protect heritage positioning while broadening reasons to engage",
        ]),
    ]
    y = 2.02
    for title, bullets in rows:
        add_arrow_panel(slide, title, bullets, 3.25, y, 11.2, 1.02)
        y += 1.18


def slide_service_white_space():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Client Service Layer & Recurring Revenue White Space",
        "Existing loyalty, bespoke, repair, and concierge touchpoints could be organized into a more repeatable care and jewelry-renewal offering",
        67,
        "Source: Sidney Garber website; Bespoke, Care & Repairs, FAQ, Concierge / Contact, Rewards references; accessed Jul-2026; diligence to confirm revenue contribution",
    )

    add_rect(slide, 1.25, 2.12, 6.25, 5.5, fill=WHITE, line=LINE)
    add_rect(slide, 8.0, 2.12, 6.25, 5.5, fill=PEACH2, line=ORANGE)
    add_text(slide, "CURRENT OFFERINGS / TOUCHPOINTS", 1.65, 2.48, 5.45, 0.28, size=13.5, bold=True, color=ORANGE, align=1)
    add_text(slide, "POTENTIAL OFFERINGS TO TEST", 8.48, 2.48, 5.28, 0.28, size=13.5, bold=True, color=ORANGE, align=1)

    # Image placeholder with camera-like abstraction, intentionally editable/no web dependency
    ph = add_rect(slide, 1.55, 3.02, 2.05, 3.9, fill=RGBColor(245, 238, 234), line=RGBColor(235, 210, 197))
    add_text(slide, "service /\nrepair\nimage", 2.0, 4.35, 1.2, 0.75, size=14, italic=True, color=GRAY, align=1)
    for i in range(5):
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.84 + i * 0.27), Inches(5.95 - i * 0.38), Inches(0.16), Inches(0.16)).fill.solid()

    offerings = [
        ("Sidney Rewards / loyalty", "Named repeat-client mechanism"),
        ("Bespoke / repurposing", "Rework stones, older pieces, and non-Sidney designs"),
        ("Care / repairs", "Repair support and ongoing care needs"),
        ("Concierge / specialist", "Direct service and appointment channel"),
    ]
    y = 3.02
    for title, detail in offerings:
        add_rect(slide, 3.78, y, 3.2, 0.72, fill=PEACH2, line=RGBColor(232, 197, 176), radius=True)
        add_text(slide, title, 4.05, y + 0.13, 2.6, 0.18, size=9.6, bold=True)
        add_text(slide, detail, 4.05, y + 0.43, 2.6, 0.14, size=6.9, color=GRAY)
        y += 0.95

    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.25), Inches(4.28), Inches(0.65), Inches(0.45))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_ORANGE
    arrow.line.fill.background()

    add_text(slide, "Structured Client Care / Jewelry Renewal Program", 8.95, 3.0, 4.4, 0.28, size=13.3, bold=True, align=1)
    tests = [
        ("Annual care plan", "Cleaning, inspection, repair coordination, sizing, refresh"),
        ("Jewelry renewal / repurpose", "Refresh inherited, unworn, or sentimental pieces into new designs"),
        ("Proactive clienteling layer", "Service reminders, occasion triggers, repeat-purchase prompts"),
    ]
    y = 3.55
    for title, detail in tests:
        add_rect(slide, 8.55, y, 5.05, 0.82, fill=WHITE, line=RGBColor(232, 197, 176), radius=True)
        add_text(slide, title, 8.90, y + 0.20, 1.65, 0.22, size=9.7, bold=True, color=ORANGE)
        add_text(slide, detail, 10.55, y + 0.19, 2.8, 0.28, size=8.4)
        y += 1.05

    add_rect(slide, 1.88, 7.82, 12.0, 0.5, fill=WHITE, line=ORANGE, radius=True)
    add_text(slide, "Diligence to confirm", 2.15, 8.0, 1.6, 0.14, size=8.8, bold=True, color=ORANGE)
    add_text(slide, "member base | purchase frequency | repair volume | service capacity | pricing / willingness to pay | margin | repeat-purchase lift",
             4.08, 8.0, 8.7, 0.14, size=8.5, color=GRAY)


def slide_growth_diligence():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Growth Diligence Priorities",
        "Each growth lever should be converted into evidence, economics, and execution ownership before it is reflected in the financial model",
        68,
    )
    x0, y0, w, h = 1.2, 2.05, 13.35, 5.95
    add_rect(slide, x0, y0, w, h, fill=WHITE, line=LINE)
    cols = [2.35, 4.25, 3.75, 3.0]
    headers = ["Lever", "Evidence to Validate", "Economics to Underwrite", "Execution Owner"]
    x = x0
    for i, cw in enumerate(cols):
        add_rect(slide, x, y0, cw, 0.55, fill=ORANGE, line=ORANGE)
        add_text(slide, headers[i], x + 0.06, y0 + 0.17, cw - 0.1, 0.16, size=10.2, bold=True, color=WHITE, align=1)
        x += cw
    rows = [
        ["Distribution expansion", "Stockist quality, sell-through, channel margin, brand fit", "Reorder rate, wholesale margin, return / markdown exposure, working capital", "Sales lead; phased account tests"],
        ["Clienteling / private client", "Customer file quality, repeat behavior, event productivity, appointment conversion", "Repeat-purchase lift, sales per event, retention, CAC by channel", "Brand / sales leader plus CRM support"],
        ["Inventory productivity", "SKU aging, cost basis, sell-through, markdown history, channel mix", "Turns, gross margin by category, liquidation value, NWC release", "Finance / merchandising cadence"],
        ["Service layer / renewal", "Repair volume, bespoke inquiries, rewards activity, service capacity", "Service margin, willingness to pay, attachment to future purchases", "Client service + operations"],
        ["Category extension", "Customer demand, brand permission, sourcing capability, competitive set", "Margin, inventory risk, required investment, test size", "Small pilots before inventory commitment"],
    ]
    row_h = (h - 0.55) / len(rows)
    y = y0 + 0.55
    for r, row in enumerate(rows):
        x = x0
        fill = PEACH2 if r % 2 == 0 else WHITE
        for c, cw in enumerate(cols):
            add_rect(slide, x, y, cw, row_h, fill=fill, line=RGBColor(225, 215, 210), width=0.5)
            add_text(slide, row[c], x + 0.12, y + 0.18, cw - 0.24, row_h - 0.2,
                     size=9.4 if c else 10.0, bold=(c == 0), color=ORANGE if c == 0 else BLACK)
            x += cw
        y += row_h
    add_text(slide, "No model credit until each lever has proof of demand, unit economics, execution owner, and working-capital impact",
             2.2, 8.12, 11.0, 0.22, size=10.8, bold=True, color=GRAY, align=1)


def slide_opportunity_case():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Preliminary Opportunity Case",
        "For Guillermo, the core question is whether attractive valuation can be paired with practical, controllable growth rather than a speculative rebrand",
        69,
    )
    boxes = [
        ("Initial Opportunity Signals", [
            "Attractive valuation",
            "Willing seller",
            "Tangible inventory / liquidation floor",
            "Kay right-to-win",
            "Wholesale / distribution upside",
        ]),
        ("What Must Be True", [
            "Brand relevance persists beyond founder / legacy story",
            "Inventory value is supported by material value, salability, and realistic markdown assumptions",
            "Channels are healthy, productive, and transferable",
            "Growth levers do not require turning the brand into mass luxury",
            "Key creative, sourcing, and client relationships can be retained or institutionalized",
        ]),
        ("Three Value Creation Levers", [
            "Commercial discipline: CRM, clienteling, events, pricing discipline",
            "Inventory productivity: turns, aging, composition, sourcing, working capital",
            "Distribution expansion: wholesale accounts, independents, private client, digital discovery",
        ]),
    ]
    xs = [1.22, 5.50, 9.78]
    widths = [3.65, 3.65, 3.65]
    for i, (title, bullets) in enumerate(boxes):
        fill = PEACH2 if i < 2 else LIGHT_GRAY
        add_rect(slide, xs[i], 2.25, widths[i], 4.95, fill=fill, line=RGBColor(220, 210, 205))
        add_text(slide, title, xs[i] + 0.25, 2.55, widths[i] - 0.5, 0.22, size=12.7, bold=True, color=BLACK, align=1)
        if i == 0:
            y = 3.1
            for idx, b in enumerate(bullets, start=1):
                add_rect(slide, xs[i] + 0.35, y, widths[i] - 0.7, 0.45, fill=LIGHT_ORANGE, line=RGBColor(240, 175, 120))
                add_text(slide, str(idx), xs[i] + 0.55, y + 0.08, 0.28, 0.2, size=18, bold=True, color=WHITE)
                add_text(slide, b, xs[i] + 1.0, y + 0.14, widths[i] - 1.35, 0.16, size=9.7)
                y += 0.68
        else:
            add_bullets(slide, bullets, xs[i] + 0.45, 3.15, widths[i] - 0.7, 3.45, size=9.3 if i == 1 else 10.0)
    add_text(slide, "Draft takeaway: attractive only if diligence confirms valuation support, seller alignment, inventory quality, and practical distribution-led growth",
             1.45, 7.55, 12.3, 0.28, size=11.4, bold=True, color=GRAY, align=1)


for builder in [slide_growth_strategy, slide_service_white_space, slide_growth_diligence, slide_opportunity_case]:
    builder()

prs.save(OUT)
print(OUT)
