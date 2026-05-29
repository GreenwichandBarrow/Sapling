#!/usr/bin/env python3
"""Build Truck Licensing & Compliance Platform one-pager from G&B template."""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from copy import deepcopy
from lxml import etree

TEMPLATE = 'brain/library/internal/one-pager-template/customs-bonds-template.pptx'
OUTPUT = '/tmp/truck-licensing-onepager.pptx'

# Section content
TITLE = "Truck Licensing & Compliance Platform — May 2026"
ASSESSMENT = "Assessment: Pending Scoring"
STATUS = "Status: Pending Scoring"

INDUSTRY_OVERVIEW = (
    "Truck Licensing & Compliance is a regulatory-driven services sector serving the "
    "trucking industry's multi-jurisdictional filing requirements. Core products: IFTA "
    "(International Fuel Tax Agreement — quarterly fuel-tax filings across 58 US/Canadian "
    "jurisdictions), IRP (International Registration Plan — apportioned vehicle "
    "registration), DOT/FMCSA authority registration, ELD compliance management, HOS "
    "(hours-of-service) auditing, driver qualification files, drug & alcohol testing "
    "consortia, and permit/title services. Buyers: small-to-mid trucking carriers "
    "(1-100 trucks), owner-operators, and fleet managers who outsource the administrative "
    "burden of multi-state compliance. Key players span scale incumbent J.J. Keller "
    "(Neenah WI, 70+ year history, $400M+ revenue), Foley Carrier Services (Hartford CT, "
    "382 employees, Greyhawk Capital-backed), DISA Global Solutions (formerly American "
    "Licensing Services — 30+ years, broad workforce-compliance roll-up), Vehicle "
    "Licensing Consultants (VLC, parent Virtual Projects LLC, founded 2000), and a long "
    "tail of ~600 NATSA-credentialed specialist firms influencing 2M+ trucks. The "
    "FMCSA's Motus registration-system launch (Q2 2026) and ELD revocation campaign "
    "(devices being purged through July 2026) are forcing carriers to seek expert help."
)

INDUSTRY_THESIS = (
    "Fragmented specialist-services market positioned as a regulatory chokepoint between "
    "trucking carriers and multi-jurisdictional regulators (FMCSA + 58 IFTA/IRP "
    "jurisdictions + state DOTs). The industry combines recurring-revenue economics "
    "(quarterly IFTA filings, annual IRP renewals, monthly ELD/HOS audit retainers) with "
    "high switching costs (state account history, DOT relationships, software-integration "
    "lift) and a meaningful long tail of ~600 NATSA-member specialist agents serving "
    "thousands of carriers — a structural pattern consistent with attractive lower-middle-"
    "market services consolidation. Mid-tier firms (10-50 FTEs) sit below the scale "
    "incumbent (J.J. Keller) and above the owner-operated single-state shops, with margin "
    "expansion available through software-enabled per-vehicle delivery."
)

DRIVERS = [
    "FMCSA enforcement escalation — Feb 2026 immediate out-of-service authority for revoked ELDs, dozens of devices purged from registered list, replacement deadlines through July 2026",
    "FMCSA Motus modernization (Q2 2026 launch) — mandatory identity verification via IDEMIA/CLEAR, stricter broker/freight-forwarder financial-responsibility rules from Jan 16, 2026",
    "Multi-state IFTA/IRP filing complexity drives outsourcing — 58 jurisdictions, quarterly cadence, audit risk on inaccurate mileage/fuel data",
    "Owner-operator outsourcing trend — first-quarter filers commonly migrate to specialists after experiencing complexity",
    "Fleet electrification creating new permitting/weight-class/state-fee categories (EV-specific registration regimes emerging)",
    "Medium-duty fleet enforcement expansion — FMCSA targeting fleets that previously assumed they were below ELD threshold",
]

RISKS = [
    "Single-source signal so far — surfaced via Helen Guo SMB Deal Hunter newsletter (2026-05-26); no second corroborating data point in this week's research window",
    "No documented G&B network access into trucking-compliance specialist community — no female-led-network anchor identified",
    "TAM and target-count estimates extrapolated, not from direct market-research reports (gated/unavailable in initial pass; deeper sourcing in Step 4)",
    "J.J. Keller scale-incumbent risk — 70+ year history, integrated software + content + services + Managed Services M&A-integration offering",
    "Fleet-tech adjacent PE consolidation (Fleetworthy + Drivewyze + Bestpass; Samsara/Geotab as platform incumbents) creates exit-channel complexity — strategic acquirers may favor telematics-integrated targets",
    "Distinct from active row 15 (Surplus Lines Compliance — state insurance tax) — different regulatory regime (FMCSA/IFTA/IRP); confirm in scoring step that the categorization stands",
]

ECONOMICS = (
    "EBITDA margins 15-25% at scale (multi-state operators with software leverage), "
    "10-15% for smaller single-jurisdiction shops. Pricing models: per-vehicle monthly "
    "subscription (typical $15-50/truck/month for IFTA+IRP+ELD bundle), quarterly retainer "
    "for filings only, transactional fees for permits/registrations, and managed-services "
    "annual contracts for full-fleet outsourcing. Revenue recurring (quarterly IFTA cycles, "
    "annual IRP renewals, ongoing audit/HOS monitoring). Retention high — switching cost "
    "is the state account history, DOT-portal credentials, and software integration. "
    "Capital-light services model. Major incumbents (J.J. Keller, Foley, DISA) supplement "
    "services with proprietary software/content. Tier sizing: J.J. Keller $400M+ "
    "consolidated revenue; mid-tier firms (Foley class) $30-100M; specialist agents in "
    "NATSA network typically $1-10M."
)

COMPETITIVE = (
    "Scale incumbent: J.J. Keller & Associates (Neenah WI, 70+ years, broadest product "
    "set — content + software + services + managed M&A integration). Mid-tier: Foley "
    "Carrier Services (Hartford CT, 382 employees, Greyhawk Capital-backed, DOT/factoring/"
    "fuel-tax/insurance), DISA Global Solutions (formerly American Licensing Services, 30+ "
    "years, broad workforce-screening roll-up acquired licensing/permitting capability), "
    "Vehicle Licensing Consultants/VLC (Virtual Projects LLC, founded 2000, IFTA/IRP/2290/"
    "DOT specialist), Transportation Compliance Service (TCS, acquired IFTA Plus 2021). "
    "Long tail: ~600 NATSA-credentialed specialist firms across US/Canada/Mexico. "
    "Adjacent: fleet-tech telematics consolidators (Samsara, Geotab, Platform Science) "
    "increasingly embedding compliance modules; Fleetworthy assembling stack via "
    "Drivewyze/Bestpass acquisitions."
)

CUSTOMERS = (
    "Small-to-mid trucking carriers (1-100 trucks) operating in 2+ IFTA/IRP jurisdictions. "
    "Owner-operators (single-truck independents who file IFTA quarterly). Fleet managers at "
    "regional/specialty carriers without in-house compliance staff. New-entrant carriers "
    "(Foley and J.J. Keller both lead with new-entrant authority packages). Medium-duty "
    "fleets newly in FMCSA enforcement scope. Buyers value: time savings, audit defense, "
    "multi-state regulatory expertise, integration with telematics/ELD providers."
)

BARRIERS = [
    "Regulatory expertise — 58 IFTA/IRP jurisdictions + FMCSA + state DOTs",
    "State-by-state account history and DOT-portal credentials",
    "ELD/permitting software integrations with major telematics platforms (Samsara, Geotab, Motive, Verizon Connect)",
    "NATSA credentialing & training network",
    "Switching costs — state account history, integrated billing, established carrier-tax-base relationships",
    "Audit-defense reputation — buyers heavily weight track record in carrier audits",
]

KSF = (
    "Multi-jurisdictional regulatory specialization (the more states covered, the larger the addressable carrier base) | "
    "Software-enabled per-vehicle service delivery (cost-to-serve drops with automation) | "
    "NATSA membership and credentialed staff (signaling + training pipeline) | "
    "Owner-operator retention through service automation (self-service portal + proactive filings) | "
    "Telematics-platform integrations (Samsara/Geotab/Motive/Verizon Connect) to capture ELD data flow | "
    "Audit-defense capability — reputational moat in a regulator-driven market"
)

EXIT = (
    "Three plausible exit channels: (1) Fleet-tech telematics consolidators — Samsara, Geotab, "
    "Platform Science, Motive — already absorbing compliance modules into broader fleet-management "
    "stacks (Fleetworthy precedent acquiring Drivewyze + Bestpass + Haul); (2) PE roll-up of "
    "pure-play compliance services — trajectory points there but not yet mature, Greyhawk's Foley "
    "and TCS/IFTA Plus the visible threads; (3) J.J. Keller as potential strategic acquirer for a "
    "tuck-in if remaining independent (its Managed Services arm already integrates acquired-fleet "
    "compliance). Caveat: adjacent fleet-tech PE consolidation may favor telematics-integrated "
    "targets, which complicates valuation for pure-services firms without software stack."
)

# Sources — grouped as required
SOURCES = {
    "Gathering-agent findings (chatroom)": [
        ("[niche-intel-recent 22:38] — brain/traces/agents/2026-05-26-niche-intelligence.md", "brain/traces/agents/2026-05-26-niche-intelligence.md"),
        ("[niche-intel-historical 22:42] — brain/traces/agents/2026-05-26-niche-intelligence.md", "brain/traces/agents/2026-05-26-niche-intelligence.md"),
        ("[niche-intel-synthesizer 22:43] — brain/traces/agents/2026-05-26-niche-intelligence.md", "brain/traces/agents/2026-05-26-niche-intelligence.md"),
        ("[niche-intel-identifier 22:47] — brain/traces/agents/2026-05-26-niche-intelligence.md", "brain/traces/agents/2026-05-26-niche-intelligence.md"),
    ],
    "External research & industry reports": [
        ("IFTA, Inc. — International Fuel Tax Association (jurisdiction list, filing schema)", "https://www.iftach.org/"),
        ("NATSA — Fleet Compliance Services overview (member network, services)", "https://mynatsa.org/fleet-compliance-services/"),
        ("NATSA — About (~600 trained specialists, 2M+ trucks influenced)", "https://mynatsa.org/about-natsa/"),
        ("J.J. Keller — Vehicle Tax & Licensing Service (IFTA/IRP product)", "https://www.jjkeller.com/shop/vehicle-tax-licensing-service"),
        ("J.J. Keller — Industry overview (70+ years compliance)", "https://www.jjkeller.com/industry/transportation-industry"),
        ("Foley — DOT Compliance services (Hartford CT, 382 employees)", "https://www.foley.io/"),
        ("Foley Carrier Services — PitchBook profile (Greyhawk Capital investor)", "https://pitchbook.com/profiles/company/144916-93"),
        ("DISA Global Solutions — DOT/Transportation Compliance", "https://disa.com/dot-transportation-compliance/"),
        ("DISA — Transportation Licensing & Permitting (formerly American Licensing Services)", "https://www.americanlicensing.com/services/trucking-interstate-licensing-permits/"),
        ("Transportation Compliance Service — PitchBook (TCS acquired IFTA Plus 2021)", "https://pitchbook.com/profiles/company/507314-89"),
        ("VLC — IFTA/IRP/2290/DOT services (Virtual Projects LLC, founded 2000)", "https://im4trux.com/"),
        ("FreightWaves — FMCSA purges dozens of ELDs amid compliance crackdown", "https://www.freightwaves.com/news/fmcsa-purges-dozens-of-elds-amid-compliance-crackdown"),
        ("FleetRabbit — ELD Mandate 2026 compliance guide", "https://fleetrabbit.com/blogs/post/eld-mandate-guide-2026"),
        ("Truck Dispatch Experts — FMCSA Rule Changes 2026 carrier guide", "https://truckdispatchexperts.com/resources/fmcsa-rules-2026/"),
        ("Heavy Duty Trucking — FMCSA Motus System launch Q2 2026", "https://www.truckinginfo.com/news/fmcsas-motus-system-is-coming-what-fleets-need-to-know-now"),
        ("Birmingham Freightliner — FMCSA 2026 Changes for fleets/drivers/brokers", "https://birminghamfreightliner.com/blog/news/fmsca-2026-changes-for-fleets-drivers-and-brokers"),
        ("Logrock — ELD Compliance 2026 rules/exemptions/penalties", "https://www.logrock.com/fmcsa-and-dot-compliance/eld-compliance/"),
        ("FreightWaves — Fleetworthy acquires Haul (AI fleet compliance roll-up precedent)", "https://www.freightwaves.com/news/fleetworthy-acquires-haul-for-ai-enhanced-fleet-compliance"),
        ("Trucking Authority — Fleet compliance during mergers/acquisitions", "https://www.truckingauthority.com/resources/articles/fleet-compliance-during-mergers-acquisitions"),
        ("PRNewswire — DISA + Tenstreet partnership (transportation compliance)", "https://www.prnewswire.com/news-releases/transforming-compliance-in-transportation-the-disa-and-tenstreet-partnership-302300244.html"),
        ("Akin — 2026 Perspectives in PE: Transportation", "https://www.akingump.com/en/insights/articles/2026-perspectives-in-private-equity-transportation"),
        ("Helen Guo — SMB Deal Hunter newsletter (single-source signal, 2026-05-26)", "https://www.smbdealhunter.com/"),
    ],
    "Internal vault references": [
        ("brain/context/learnings.md — niche-intelligence prior runs", "brain/context/learnings.md"),
        (".claude/skills/niche-intelligence/references/one-pager-template.md — template spec", ".claude/skills/niche-intelligence/references/one-pager-template.md"),
    ],
    "CRM / data pulls": [
        ("N/A this run — no Attio/Apollo/Linkt query performed; deeper sourcing in Step 4 scorer", ""),
    ],
}


def set_cell_text(cell, text, font_size_pt=10, bold=False, color_rgb=None):
    """Wipe cell and write text into a single new paragraph/run with controlled formatting."""
    tf = cell.text_frame
    # Clear all paragraphs except keep one
    # Remove all <a:p> children, then add fresh
    txBody = tf._txBody
    for p in txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        txBody.remove(p)
    # Add first paragraph
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 or len(tf.paragraphs) == 0 else tf.paragraphs[0]
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size_pt)
        if bold:
            run.font.bold = True
        if color_rgb is not None:
            run.font.color.rgb = color_rgb


def set_title_cell(cell, text):
    """Special-cased title cell — set explicit black color + 16pt per template reference."""
    tf = cell.text_frame
    txBody = tf._txBody
    for p in txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        txBody.remove(p)
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)


def set_header_cell(cell, text):
    """Section header — bold 14pt."""
    tf = cell.text_frame
    txBody = tf._txBody
    for p in txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        txBody.remove(p)
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True


def set_body_cell(cell, text):
    """Body cell — 9pt non-bold."""
    set_cell_text(cell, text, font_size_pt=9, bold=False)


def set_bullet_cell(cell, items):
    """List of items, one per line, 9pt."""
    text = '\n'.join(items)
    set_cell_text(cell, text, font_size_pt=9, bold=False)


def main():
    prs = Presentation(TEMPLATE)
    slide = prs.slides[0]
    tbl = slide.shapes[2].table

    # Row 0: Title (merge cells in template — col 0 has the text)
    set_title_cell(tbl.cell(0, 0), TITLE)

    # Row 1: Assessment / Status
    set_cell_text(tbl.cell(1, 0), ASSESSMENT, font_size_pt=12, bold=True)
    set_cell_text(tbl.cell(1, 1), STATUS, font_size_pt=12, bold=True)

    # Row 2: Industry Overview header
    set_header_cell(tbl.cell(2, 0), "Industry Overview")

    # Row 3: Industry Overview body
    set_body_cell(tbl.cell(3, 0), INDUSTRY_OVERVIEW)

    # Row 4: Industry Thesis header
    set_header_cell(tbl.cell(4, 0), "Industry Thesis")

    # Row 5: Industry Thesis body
    set_body_cell(tbl.cell(5, 0), INDUSTRY_THESIS)

    # Row 6: Drivers | Risks headers
    set_header_cell(tbl.cell(6, 0), "Macro Trends & Growth Drivers")
    set_header_cell(tbl.cell(6, 1), "Risks & Concerns")

    # Row 7: Drivers | Risks body
    set_bullet_cell(tbl.cell(7, 0), DRIVERS)
    set_bullet_cell(tbl.cell(7, 1), RISKS)

    # Row 8: Economics | Competitive headers
    set_header_cell(tbl.cell(8, 0), "Economics & Pricing")
    set_header_cell(tbl.cell(8, 1), "Competitive Landscape")

    # Row 9: Economics | Competitive body
    set_body_cell(tbl.cell(9, 0), ECONOMICS)
    set_body_cell(tbl.cell(9, 1), COMPETITIVE)

    # Row 10: Customers | Barriers headers
    set_header_cell(tbl.cell(10, 0), "Customers")
    set_header_cell(tbl.cell(10, 1), "Barriers to Entry")

    # Row 11: Customers | Barriers body
    set_body_cell(tbl.cell(11, 0), CUSTOMERS)
    set_bullet_cell(tbl.cell(11, 1), BARRIERS)

    # Row 12: Key Success Factors header
    set_header_cell(tbl.cell(12, 0), "Key Success Factors")

    # Row 13: Key Success Factors body (template had split — write KSF into col 0, blank col 1)
    set_body_cell(tbl.cell(13, 0), KSF)
    set_body_cell(tbl.cell(13, 1), "")

    # Row 14: Exit header
    set_header_cell(tbl.cell(14, 0), "Exit")

    # Row 15: Exit body
    set_body_cell(tbl.cell(15, 0), EXIT)

    # SOURCES slide — add a second slide
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    sources_slide = prs.slides.add_slide(blank_layout)

    # Add title text box
    from pptx.util import Inches
    title_box = sources_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.6))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = "Sources — Truck Licensing & Compliance Platform"
    title_run.font.size = Pt(20)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 0, 0)

    # Add body text box with all sources grouped
    body_box = sources_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.5), Inches(6.5))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True

    first_para = True
    for group_name, entries in SOURCES.items():
        # Group header
        p = body_tf.paragraphs[0] if first_para else body_tf.add_paragraph()
        first_para = False
        run = p.add_run()
        run.text = group_name
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)

        # Entries
        for title_text, url in entries:
            p2 = body_tf.add_paragraph()
            r = p2.add_run()
            r.text = f"• {title_text}"
            r.font.size = Pt(8)
            if url:
                r.hyperlink.address = url
                r.font.color.rgb = RGBColor(0x0B, 0x5C, 0xAB)
            else:
                r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Spacer
        sp = body_tf.add_paragraph()
        sp_run = sp.add_run()
        sp_run.text = " "
        sp_run.font.size = Pt(4)

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")

    # Count sources
    total = sum(len(v) for v in SOURCES.values())
    print(f"Sources count: {total} entries across {len(SOURCES)} groups")


if __name__ == '__main__':
    main()
